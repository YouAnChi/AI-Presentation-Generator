"""排版构建（Builder）智能体：将结构化的演示内容写入 `.pptx` 文件。

核心职责：
- 根据 `PresentationOutline` 的每页 `SlideContent`，选择合适布局，填充标题、正文与备注；
- 按需插入图片（若存在有效路径），并做简单的版式放置；
- 输出最终文件到 `output/` 目录，返回保存路径。

实现要点：
- 使用 `python-pptx` 操作模板与形状；
- 简化的布局索引映射（Title/Title and Content/Section Header/Two Content）；
- 尝试不同插图位置，并进行异常捕获，避免中断整体生成。
"""

import os
import json
import logging
from typing import AsyncIterable, Any
from pptx import Presentation
from ai_ppt.common.types import PresentationOutline
from ai_ppt.common.utils import get_logger
from ai_ppt.common.base_agent import BaseAgent
from a2a.types import AgentCard

logger = get_logger(__name__)

class BuilderAgent(BaseAgent):
    """PPT 构建智能体：将结构化内容落盘为 `.pptx`。"""
    def __init__(self, output_dir: str = "output"):
        super().__init__(
            agent_name="PPT Builder",
            description="Compiles content into a final PowerPoint file."
        )
        self.output_dir = output_dir
        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

    def build_presentation(self, outline: PresentationOutline, filename: str = "presentation.pptx") -> str:
        """根据结构化大纲构建演示文稿并保存到本地。

        参数：
            outline: 演示文稿结构（主题、每页内容、版式）。
            filename: 输出文件名，默认 `presentation.pptx`。

        返回：
            保存后的绝对路径或相对路径字符串。
        """
        logger.info(f"Building presentation: {outline.topic}")
        prs = Presentation()
        
        for slide_content in outline.slides:
            # 简单布局映射：根据字符串包含关系选择模板索引
            if "Title Slide" in slide_content.layout:
                layout_index = 0
            elif "Title and Content" in slide_content.layout:
                layout_index = 1
            elif "Section Header" in slide_content.layout:
                layout_index = 2
            elif "Two Content" in slide_content.layout:
                layout_index = 3
            else:
                layout_index = 1 
            
            slide_layout = prs.slide_layouts[layout_index]
            slide = prs.slides.add_slide(slide_layout)
            
            # 写入标题
            if slide.shapes.title:
                slide.shapes.title.text = slide_content.title
            
            # 写入正文（占位符索引 1 通常为正文区域）
            if len(slide.placeholders) > 1 and slide_content.body_text:
                body_shape = slide.placeholders[1]
                if hasattr(body_shape, "text"):
                    body_shape.text = slide_content.body_text
            
            # 写入演讲者备注（notes 页面）
            if slide_content.speaker_notes and slide.has_notes_slide:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_content.speaker_notes

            # 插入图片（若路径存在）
            if slide_content.image_path and os.path.exists(slide_content.image_path):
                try:
                    # 简化放置策略：不同布局选择不同放置位置与尺寸
                    left = top = 0
                    if layout_index == 1: # Title and Content
                        # 将图片放在右侧，并缩放宽度
                        left = prs.slide_width * 0.6
                        top = prs.slide_height * 0.3
                        width = prs.slide_width * 0.35
                        slide.shapes.add_picture(slide_content.image_path, left, top, width=width)
                    elif layout_index == 3: # Two Content
                         # 若存在第二个占位符，插入其中
                         if len(slide.placeholders) > 2:
                             p = slide.placeholders[2]
                             p.insert_picture(slide_content.image_path)
                    else:
                        # 默认放置：靠左下方，按高度缩放
                        left = prs.slide_width * 0.1
                        top = prs.slide_height * 0.5
                        height = prs.slide_height * 0.4
                        slide.shapes.add_picture(slide_content.image_path, left, top, height=height)
                        
                except Exception as e:
                    logger.error(f"Failed to add image to slide {slide_content.title}: {e}")

        output_path = os.path.join(self.output_dir, filename)
        prs.save(output_path)
        logger.info(f"Presentation saved to: {output_path}")
        return output_path

    async def stream(
        self, query: str, context_id: str, task_id: str
    ) -> AsyncIterable[dict[str, Any]]:
        """流式构建 `.pptx` 文件：接收结构化大纲，生成文件并返回路径。"""
        logger.info(f"Received request to build PPT")
        
        try:
            # 预期输入为 `PresentationOutline` 的 JSON 字符串
            data = json.loads(query)
            outline = PresentationOutline(**data)
            
            output_path = self.build_presentation(outline)
            
            yield self.format_response(f"Presentation built successfully: {output_path}")
            
        except Exception as e:
            logger.error(f"Error building PPT: {e}")
            yield self.format_response(f"Error: {e}")

    def get_agent_card(self, host: str, port: int) -> AgentCard:
        """生成当前智能体的 AgentCard，并声明其工具与标签。"""
        card = super().get_agent_card(host, port)
        card.skills = [
            {
                "id": "build_ppt",
                "name": "Build PPT File",
                "description": "Generates a .pptx file from structured slide data",
                "tags": ["builder", "pptx"]
            }
        ]
        return card

if __name__ == "__main__":
    import click
    from ai_ppt.common.server_utils import start_agent_server
    
    @click.command()
    @click.option("--host", default="localhost")
    @click.option("--port", default=10203)
    def main(host, port):
        # 启动 Builder 的 A2A 服务端
        agent = BuilderAgent()
        start_agent_server(agent, host, port)
        
    main()
